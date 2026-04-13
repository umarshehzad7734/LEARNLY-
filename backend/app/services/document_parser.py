"""
Enhanced Document Parser Module
Integrates RAGFlow's advanced document parsing capabilities for better text extraction
"""

import os
import re
from typing import List, Dict, Tuple, Optional
from pathlib import Path
import base64
from io import BytesIO

# PDF parsing
from pypdf import PdfReader
import pdfplumber

# DOCX parsing
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph

# Markdown parsing
import markdown
from bs4 import BeautifulSoup

# Excel parsing
import openpyxl

# PPTX parsing
from pptx import Presentation

# Image handling
from PIL import Image


class DocumentParser:
    """
    Advanced document parser that extracts text, tables, and images from various file formats
    """
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.txt', '.md', '.xlsx', '.xls', '.pptx', '.ppt']
    
    def parse(self, file_path: str) -> Dict:
        """
        Parse a document and extract all content
        
        Returns:
            Dict with keys: text, tables, images, metadata
        """
        file_extension = Path(file_path).suffix.lower()
        
        if file_extension == '.pdf':
            return self.parse_pdf(file_path)
        elif file_extension == '.docx':
            return self.parse_docx(file_path)
        elif file_extension == '.txt':
            return self.parse_txt(file_path)
        elif file_extension == '.md':
            return self.parse_markdown(file_path)
        elif file_extension in ['.xlsx', '.xls']:
            return self.parse_excel(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            return self.parse_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    
    def parse_pdf(self, file_path: str) -> Dict:
        """
        Enhanced PDF parsing: Fast text extraction with pypdf, 
        and targeted table extraction with pdfplumber.
        """
        result = {
            'text': '',
            'tables': [],
            'images': [],
            'metadata': {
                'pages': 0,
                'has_tables': False,
                'has_images': False
            }
        }
        
        try:
            # 1. Fast text extraction using pypdf
            reader = PdfReader(file_path)
            result['metadata']['pages'] = len(reader.pages)
            
            pypdf_text = ""
            for page_num, page in enumerate(reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    pypdf_text += f"\n\n--- Page {page_num} ---\n\n{page_text}"
                
                # Check for images
                if '/XObject' in page['/Resources']:
                    xObject = page['/Resources']['/XObject'].get_object()
                    for obj in xObject:
                        if xObject[obj]['/Subtype'] == '/Image':
                            result['metadata']['has_images'] = True
                            result['images'].append({
                                'page': page_num,
                                'name': obj
                            })
            
            # 2. Targeted table extraction using pdfplumber (the slow part)
            # We only do this if we actually need tables, to save time
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    tables = page.extract_tables()
                    if tables:
                        result['metadata']['has_tables'] = True
                        for table_idx, table in enumerate(tables):
                            if table:
                                table_md = self._table_to_markdown(table)
                                result['tables'].append({
                                    'page': page_num,
                                    'index': table_idx,
                                    'content': table_md,
                                    'raw': table
                                })
                                # Insert table markers in text (optional, but good for context)
                                pypdf_text += f"\n\n[Table {table_idx + 1} on Page {page_num}]\n{table_md}\n"
            
            result['text'] = pypdf_text
            
        except Exception as e:
            print(f"Error parsing PDF {file_path}: {str(e)}")
            # Ultra-basic fallback
            try:
                reader = PdfReader(file_path)
                result['text'] = "".join([p.extract_text() or "" for p in reader.pages])
                result['metadata']['pages'] = len(reader.pages)
            except Exception as e2:
                print(f"Fallback PDF parsing also failed: {str(e2)}")
                raise
        
        return result
        
        return result
    
    def parse_docx(self, file_path: str) -> Dict:
        """
        Enhanced DOCX parsing with tables and images
        """
        result = {
            'text': '',
            'tables': [],
            'images': [],
            'metadata': {
                'paragraphs': 0,
                'has_tables': False,
                'has_images': False
            }
        }
        
        try:
            doc = Document(file_path)
            
            # Extract content in document order (paragraphs and tables)
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    # Paragraph
                    paragraph = Paragraph(element, doc)
                    text = paragraph.text.strip()
                    if text:
                        result['text'] += text + "\n"
                        result['metadata']['paragraphs'] += 1
                    
                    # Check for images in paragraph
                    for run in paragraph.runs:
                        if run._element.xpath('.//pic:pic'):
                            result['metadata']['has_images'] = True
                            result['images'].append({
                                'paragraph': result['metadata']['paragraphs'],
                                'type': 'inline'
                            })
                
                elif isinstance(element, CT_Tbl):
                    # Table
                    table = Table(element, doc)
                    result['metadata']['has_tables'] = True
                    
                    # Extract table data
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    # Convert to markdown
                    table_md = self._table_to_markdown(table_data)
                    result['tables'].append({
                        'index': len(result['tables']),
                        'content': table_md,
                        'raw': table_data
                    })
                    
                    # Add table to text
                    result['text'] += f"\n\n[Table {len(result['tables'])}]\n{table_md}\n\n"
        
        except Exception as e:
            print(f"Error parsing DOCX {file_path}: {str(e)}")
            raise
        
        return result
    
    def parse_txt(self, file_path: str) -> Dict:
        """
        Parse plain text file
        """
        result = {
            'text': '',
            'tables': [],
            'images': [],
            'metadata': {}
        }
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                result['text'] = f.read()
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        result['text'] = f.read()
                    break
                except:
                    continue
        
        return result
    
    def parse_markdown(self, file_path: str) -> Dict:
        """
        Parse Markdown file
        """
        result = {
            'text': '',
            'tables': [],
            'images': [],
            'metadata': {
                'has_code_blocks': False,
                'has_tables': False
            }
        }
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Convert markdown to HTML for parsing
            html = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
            soup = BeautifulSoup(html, 'html.parser')
            
            # Extract text
            result['text'] = soup.get_text()
            
            # Check for code blocks
            if '```' in md_content or '<code>' in html:
                result['metadata']['has_code_blocks'] = True
            
            # Extract tables
            tables = soup.find_all('table')
            if tables:
                result['metadata']['has_tables'] = True
                for idx, table in enumerate(tables):
                    rows = []
                    for tr in table.find_all('tr'):
                        row = [td.get_text().strip() for td in tr.find_all(['td', 'th'])]
                        rows.append(row)
                    
                    table_md = self._table_to_markdown(rows)
                    result['tables'].append({
                        'index': idx,
                        'content': table_md,
                        'raw': rows
                    })
        
        except Exception as e:
            print(f"Error parsing Markdown {file_path}: {str(e)}")
            raise
        
        return result
    
    def parse_excel(self, file_path: str) -> Dict:
        """
        Parse Excel file
        """
        result = {
            'text': '',
            'tables': [],
            'images': [],
            'metadata': {
                'sheets': []
            }
        }
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                result['metadata']['sheets'].append(sheet_name)
                
                # Extract sheet data
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None and str(cell).strip() for cell in row):
                        sheet_data.append([str(cell) if cell is not None else '' for cell in row])
                
                if sheet_data:
                    # Convert to markdown
                    table_md = self._table_to_markdown(sheet_data)
                    result['tables'].append({
                        'sheet': sheet_name,
                        'content': table_md,
                        'raw': sheet_data
                    })
                    
                    # Add to text
                    result['text'] += f"\n\n--- Sheet: {sheet_name} ---\n\n{table_md}\n\n"
        
        except Exception as e:
            print(f"Error parsing Excel {file_path}: {str(e)}")
            raise
        
        return result
    
    def parse_pptx(self, file_path: str) -> Dict:
        """
        Parse PowerPoint file
        """
        result = {
            'text': '',
            'tables': [],
            'images': [],
            'metadata': {
                'slides': 0,
                'has_tables': False,
                'has_images': False
            }
        }
        
        try:
            prs = Presentation(file_path)
            result['metadata']['slides'] = len(prs.slides)
            
            full_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                # Use a slide header for context
                slide_text.append(f"--- Slide {i} ---")
                
                for shape in slide.shapes:
                    # Extract text from shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract tables from shapes
                    if shape.has_table:
                        result['metadata']['has_tables'] = True
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text_frame.text.strip() for cell in row.cells]
                            table_data.append(row_data)
                        
                        if table_data:
                            table_md = self._table_to_markdown(table_data)
                            result['tables'].append({
                                'slide': i,
                                'content': table_md,
                                'raw': table_data
                            })
                            slide_text.append(f"\n[Table on Slide {i}]\n{table_md}\n")
                    
                    # Check for images (Picture shapes)
                    # Note: We don't extract the actual image bytes here yet, 
                    # just mark that they exist
                    if hasattr(shape, "shape_type"):
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            result['metadata']['has_images'] = True
                            result['images'].append({
                                'slide': i,
                                'name': getattr(shape, 'name', f"Picture {len(result['images']) + 1}")
                            })

                full_text.append("\n".join(slide_text))
            
            result['text'] = "\n\n".join(full_text)
            
        except Exception as e:
            print(f"Error parsing PPTX {file_path}: {str(e)}")
            raise
        
        return result
    
    def _table_to_markdown(self, table_data: List[List[str]]) -> str:
        """
        Convert table data to markdown format
        """
        if not table_data:
            return ""
        
        # Clean and prepare data
        cleaned_data = []
        for row in table_data:
            cleaned_row = [str(cell).strip() if cell else '' for cell in row]
            cleaned_data.append(cleaned_row)
        
        if not cleaned_data:
            return ""
        
        # Determine column widths
        max_cols = max(len(row) for row in cleaned_data)
        col_widths = [0] * max_cols
        
        for row in cleaned_data:
            for i, cell in enumerate(row):
                if i < max_cols:
                    col_widths[i] = max(col_widths[i], len(cell))
        
        # Build markdown table
        md_lines = []
        
        for idx, row in enumerate(cleaned_data):
            # Pad row to max columns
            padded_row = row + [''] * (max_cols - len(row))
            
            # Format cells
            formatted_cells = [cell.ljust(col_widths[i]) for i, cell in enumerate(padded_row)]
            md_lines.append('| ' + ' | '.join(formatted_cells) + ' |')
            
            # Add separator after first row (header)
            if idx == 0:
                separators = ['-' * col_widths[i] for i in range(max_cols)]
                md_lines.append('| ' + ' | '.join(separators) + ' |')
        
        return '\n'.join(md_lines)


# Singleton instance
document_parser = DocumentParser()
